"""PPTX presentation ingestion using python-pptx."""

from __future__ import annotations

from pathlib import Path

from ingestion.base import Document, Ingestor


class PptxIngestor(Ingestor):
    """Extract slide text, speaker notes, and images from PPTX files."""

    def ingest(self, file_path: str | Path) -> Document:
        path = self._validate_path(file_path)
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError("python-pptx is required: pip install python-pptx") from exc

        text_chunks: list[str] = []
        images: list[bytes] = []

        prs = Presentation(str(path))
        metadata = {
            "format": "pptx",
            "filename": path.name,
            "slides": str(len(prs.slides)),
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text.strip())

                if shape.shape_type == 13:  # Picture type
                    try:
                        images.append(shape.image.blob)
                    except Exception:
                        continue

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes] {notes}")

            if len(slide_texts) > 1:
                text_chunks.append("\n".join(slide_texts))

        return Document(text_chunks=text_chunks, images=images, metadata=metadata)
